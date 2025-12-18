import os
import sys
from datetime import datetime
from typing import Dict, Any, List

# Try importing pptx, with dynamic pip fallback
try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    try:
        import subprocess
        # Proactively install python-pptx
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        HAS_PPTX = True
    except Exception:
        HAS_PPTX = False

class SlidesExporterModule:
    """
    Generates PowerPoint slide presentations (.pptx) based on dataset analysis.
    """
    def __init__(self):
        pass

    def build_pptx_report(
        self,
        analysis_data: Dict[str, Any],
        output_path: str,
        target_column: str = None
    ) -> bool:
        """
        Creates a multi-slide presentation summarizing the dataset audit and ML results.
        """
        if not HAS_PPTX:
            raise RuntimeError("The python-pptx library is missing and could not be dynamically installed.")

        prs = Presentation()
        
        # Dark Slate Theme Colors
        BG_COLOR = RGBColor(15, 23, 42)       # Slate 900
        TEXT_COLOR = RGBColor(241, 245, 249)  # Slate 100
        MUTED_COLOR = RGBColor(148, 163, 184) # Slate 400
        ACCENT_COLOR = RGBColor(99, 102, 241) # Indigo 500
        PASS_COLOR = RGBColor(16, 185, 129)   # Emerald 500

        # Helper: Set dark slide background
        def set_slide_background(slide):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = BG_COLOR

        # Helper: Create title on slide
        def add_slide_header(slide, title_text, category_text=None):
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(1.0))
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.margin_top = Inches(0)
            tf.margin_left = Inches(0)
            
            if category_text:
                p_cat = tf.paragraphs[0]
                p_cat.text = category_text.upper()
                p_cat.font.size = Pt(11)
                p_cat.font.bold = True
                p_cat.font.color.rgb = ACCENT_COLOR
                p_cat.font.name = "Arial"
                p_cat.space_after = Pt(2)
                p_title = tf.add_paragraph()
            else:
                p_title = tf.paragraphs[0]
                
            p_title.text = title_text
            p_title.font.size = Pt(28)
            p_title.font.bold = True
            p_title.font.color.rgb = TEXT_COLOR
            p_title.font.name = "Arial"

        # --- SLIDE 1: Title Slide (Layout 6 is blank) ---
        blank_slide_layout = prs.slide_layouts[6]
        slide1 = prs.slides.add_slide(blank_slide_layout)
        set_slide_background(slide1)

        # Draw Title Card
        title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(2.5))
        tf1 = title_box.text_frame
        tf1.word_wrap = True
        
        p1 = tf1.paragraphs[0]
        p1.text = "AURAEDA V3.0 EXECUTIVE AUDIT"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_COLOR
        p1.font.name = "Arial"
        p1.space_after = Pt(8)

        p2 = tf1.add_paragraph()
        p2.text = "Dataset Health & Model Prep Report"
        p2.font.size = Pt(36)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_COLOR
        p2.font.name = "Arial"
        p2.space_after = Pt(12)

        p3 = tf1.add_paragraph()
        curr_time = datetime.now().strftime('%B %d, %Y')
        p3.text = f"Generated Automatically on {curr_time}\nInteractive ML Preprocessing Sandbox Suite"
        p3.font.size = Pt(13)
        p3.font.color.rgb = MUTED_COLOR
        p3.font.name = "Arial"

        # --- SLIDE 2: Executive Summary ---
        slide2 = prs.slides.add_slide(blank_slide_layout)
        set_slide_background(slide2)
        add_slide_header(slide2, "Dataset Integrity Score & Summary", "Executive Overview")

        # Add metric stats boxes (integrity score + rows/columns)
        summary = analysis_data.get("dataset_summary", {})
        score = analysis_data.get("results", {}).get("alerts", {}).get("health_score", 100)
        
        box1 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(4.5))
        tf2 = box1.text_frame
        tf2.word_wrap = True
        
        p = tf2.paragraphs[0]
        p.text = "DATASET AUDIT SIGNALS"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        p.space_after = Pt(10)

        bullets = [
            f"Integrity Score: {score}/100",
            f"Total Instances: {summary.get('n_rows', 0):,} rows",
            f"Total Variables: {summary.get('n_columns', 0)} columns",
            f"Missing Values: {summary.get('null_percentage', 0.0):.2f}% of total cells",
            f"Memory Footprint: {summary.get('memory_usage_mb', 0.0):.3f} MB",
            f"Target Column: '{target_column}'" if target_column else "Target Column: [Not Selected]"
        ]
        for b in bullets:
            p = tf2.add_paragraph()
            p.text = "• " + b
            p.font.size = Pt(15)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(6)

        # Right side box: Commentary bullet points
        box2 = slide2.shapes.add_textbox(Inches(5.0), Inches(1.8), Inches(4.5), Inches(4.5))
        tf3 = box2.text_frame
        tf3.word_wrap = True
        
        p = tf3.paragraphs[0]
        p.text = "KEY OBSERVATIONS & DOWNSTREAM IMPACT"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        p.space_after = Pt(10)

        # Add smart observations
        obs = [
            "Low health score indicates potential structural errors (duplicates, PII leakage, etc.) requiring wrangling.",
            "Variables with high missingness must be dropped or imputed prior to modeling to avoid sample bias.",
            "Validate target classes for severe skew; consider class weights or SMOTE balancing in the model sandbox."
        ]
        for o in obs:
            p = tf3.add_paragraph()
            p.text = "▪ " + o
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(10)

        # --- SLIDE 3: Quality Alerts Catalog ---
        slide3 = prs.slides.add_slide(blank_slide_layout)
        set_slide_background(slide3)
        add_slide_header(slide3, "Data Anomalies & Alerts Catalog", "Diagnose & Audit")

        alerts_list = analysis_data.get("results", {}).get("alerts", {}).get("alerts", [])
        if not alerts_list:
            alerts_list = [{"type": "Info", "risk": "low", "message": "No major data quality alerts were found."}]

        # Create a table for alerts
        rows = min(len(alerts_list) + 1, 6) # show up to 5 alerts
        cols = 3
        
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9.0)
        height = Inches(0.5 * rows)

        table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Column Widths
        table.columns[0].width = Inches(1.5) # Type
        table.columns[1].width = Inches(1.2) # Risk
        table.columns[2].width = Inches(6.3) # Message

        # Table Header
        headers = ["Anomaly Type", "Risk Level", "Audit Message"]
        for col_idx, h in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_COLOR
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = TEXT_COLOR
                p.alignment = PP_ALIGN.LEFT

        # Fill table rows
        for row_idx in range(1, rows):
            alert = alerts_list[row_idx - 1]
            risk = str(alert.get("risk", "low")).upper()
            msg = alert.get("message", "")
            al_type = alert.get("type", "General")

            vals = [al_type, risk, msg]
            for col_idx, v in enumerate(vals):
                cell = table.cell(row_idx, col_idx)
                cell.text = v
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate 800
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.color.rgb = TEXT_COLOR
                    if col_idx == 1:
                        # Highlight risk color
                        if risk == "HIGH":
                            p.font.color.rgb = RGBColor(239, 68, 68) # Red
                        elif risk == "MEDIUM":
                            p.font.color.rgb = RGBColor(245, 158, 11) # Orange
                        else:
                            p.font.color.rgb = PASS_COLOR
                    p.alignment = PP_ALIGN.LEFT

        # --- SLIDE 4: Feature Importance & PCA ---
        slide4 = prs.slides.add_slide(blank_slide_layout)
        set_slide_background(slide4)
        add_slide_header(slide4, "Feature Relevance & Projections", "Explore & Model Prep")

        importance_data = analysis_data.get("results", {}).get("importance", {})
        top_feats = importance_data.get("mutual_info", [])[:5] if importance_data else []
        
        box1 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(4.5))
        tf4 = box1.text_frame
        tf4.word_wrap = True
        
        p = tf4.paragraphs[0]
        p.text = "MUTUAL INFORMATION IMPORTANCE"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        p.space_after = Pt(10)

        if not top_feats:
            p = tf4.add_paragraph()
            p.text = "No importance rankings calculated. Please select a target variable first."
            p.font.size = Pt(13)
            p.font.color.rgb = MUTED_COLOR
        else:
            for idx, feat in enumerate(top_feats):
                col_name = feat.get("column", "Unknown")
                score_val = feat.get("score", 0.0)
                p = tf4.add_paragraph()
                p.text = f"{idx+1}. {col_name} (Score: {score_val:.4f})"
                p.font.size = Pt(14)
                p.font.color.rgb = TEXT_COLOR
                p.space_after = Pt(8)

        # Right box: Dimensionality reduction (PCA)
        box2 = slide4.shapes.add_textbox(Inches(5.0), Inches(1.8), Inches(4.5), Inches(4.5))
        tf5 = box2.text_frame
        tf5.word_wrap = True
        
        p = tf5.paragraphs[0]
        p.text = "DIMENSIONALITY COMPRESSION (PCA)"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        p.space_after = Pt(10)

        pca_summary = [
            "Principal Component Analysis aggregates high-dimensional column variance.",
            "Top components summarize latent directions of variance in numerical fields.",
            "PCA projection is mapped into 2D and 3D scatter plots in the frontend Importance & PCA tab."
        ]
        for pt in pca_summary:
            p = tf5.add_paragraph()
            p.text = "▪ " + pt
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(10)

        # --- SLIDE 5: AutoML Leaderboard & Diagnostics ---
        slide5 = prs.slides.add_slide(blank_slide_layout)
        set_slide_background(slide5)
        add_slide_header(slide5, "AutoML Model Sandbox Benchmark", "Model Prep & AutoML")

        # Add note about leaderboards
        box_auto = slide5.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9.0), Inches(4.5))
        tf_auto = box_auto.text_frame
        tf_auto.word_wrap = True
        
        p = tf_auto.paragraphs[0]
        p.text = "AUTOML LEADERBOARD AUDITING"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        p.space_after = Pt(10)

        auto_bullets = [
            "Concurrently trains 8 estimators (Ridge/Logistic, Decision Tree, Random Forest, Naive Bayes, KNN, Boosting).",
            "Calculates metrics (F1-macro, Accuracy, AUC, R-squared) across training and testing data splits.",
            "Integrates live classification probability cutoff sliders to re-calculate confusion matrices on the fly.",
            "Fits Cook's distance diagnostics analytically to isolate highly influential outliers on regression tasks."
        ]
        for bullet in auto_bullets:
            p = tf_auto.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(15)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(8)

        # Save presentation
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        return True
